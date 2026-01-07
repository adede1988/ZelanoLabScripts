# -*- coding: utf-8 -*-
"""
Created on Thu Nov 13 23:18:16 2025

@author: dtf8829
"""

import os
import glob
import random

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# ============================================================
# USER CONFIG
# ============================================================

rootPath = r"R:\Neurology\Zelano_Lab\Lab_Common\Adam\Dupi_processing"  # <-- CHANGE THIS

sessionIDs = [
    "251009_OBE_NWU_CP_1",  # CP
    "250225_OBE_NWU_AS_4",  # AS
    "250904_OBE_NWU_TI",    # TI
]

# Map full session IDs to short labels and order
subj_short = {
    "251009_OBE_NWU_CP_1": "CP",
    "250225_OBE_NWU_AS_4": "AS",
    "250904_OBE_NWU_TI":   "TI",
}

# Order of subjects (columns)
subject_order = ["251009_OBE_NWU_CP_1",
                 "250225_OBE_NWU_AS_4",
                 "250904_OBE_NWU_TI"]

# Output file
output_pptx = "EmotionalMovieTask_MultiSubject.pptx"

# Make random choices reproducible for the *evt****.jpg files
random.seed(1)

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def fig_path_for_pattern(rootPath, sessionID, pattern):
    """
    Return the full path to an image matching pattern.
    If pattern contains '*', use glob and randomly pick one match.
    If pattern is exact, just return it if it exists.
    Returns None if no file is found.
    """
    task_dir = os.path.join(rootPath, sessionID, "EmotionalMovieTask")

    # Wildcard case
    if "*" in pattern:
        matches = glob.glob(os.path.join(task_dir, pattern))
        if not matches:
            print(f"[WARN] No matches for pattern '{pattern}' in {task_dir}")
            return None
        return random.choice(matches)

    # Exact filename case
    full = os.path.join(task_dir, pattern)
    if not os.path.exists(full):
        print(f"[WARN] File not found: {full}")
        return None
    return full


def add_slide_with_grid(prs, slide_title, row_patterns, rootPath, subject_order, subj_short):
    """
    Add a slide with:
      - Big title at top
      - 3 column headers (CP, AS, TI)
      - For each row_patterns entry, one image per subject (3 columns)
    row_patterns: list of filename patterns (one per row, same patterns across subjects)
    """

    # Blank layout
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Slide size and layout constants
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Convert EMU to inches helper
    def emu_to_in(x):
        return x / 914400.0

    # Margins
    top_margin_in = 0.5
    left_margin_in = 0.5
    right_margin_in = 0.5
    col_spacing_in = 0.2
    row_spacing_in = 0.2

    # Title box
    title_height_in = 0.6
    title_box = slide.shapes.add_textbox(
        Inches(left_margin_in),
        Inches(top_margin_in),
        slide_width - Inches(left_margin_in + right_margin_in),
        Inches(title_height_in),
    )
    title_tf = title_box.text_frame
    title_tf.text = slide_title
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Column header row
    header_height_in = 0.4
    n_cols = len(subject_order)

    usable_width_in = emu_to_in(slide_width) - left_margin_in - right_margin_in - (n_cols - 1) * col_spacing_in
    col_width_in = usable_width_in / n_cols

    header_top_in = top_margin_in + title_height_in + 0.1

    for col_idx, sid in enumerate(subject_order):
        short_name = subj_short.get(sid, sid)
        left_in = left_margin_in + col_idx * (col_width_in + col_spacing_in)
        tb = slide.shapes.add_textbox(
            Inches(left_in),
            Inches(header_top_in),
            Inches(col_width_in),
            Inches(header_height_in),
        )
        tf = tb.text_frame
        tf.text = short_name
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Image rows
    n_rows = len(row_patterns)
    # Rough guess for picture height to fit all rows:
    remaining_height_in = emu_to_in(slide_height) - (header_top_in + header_height_in + 0.2) - 0.5
    pic_height_in = (remaining_height_in - (n_rows - 1) * row_spacing_in) / max(n_rows, 1)

    current_top_in = header_top_in + header_height_in + 0.2

    for row_idx, pattern in enumerate(row_patterns):
        for col_idx, sid in enumerate(subject_order):
            img_path = fig_path_for_pattern(rootPath, sid, pattern)
            if img_path is None:
                continue  # leave blank if missing

            left_in = left_margin_in + col_width_in * col_idx + col_spacing_in * col_idx
            slide.shapes.add_picture(
                img_path,
                Inches(left_in),
                Inches(current_top_in),
                height=Inches(pic_height_in),
            )

        current_top_in += pic_height_in + row_spacing_in


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation()

# ----------------------------------------------------------------
# Slide 1: Gamma peak frequency
# ----------------------------------------------------------------
slide1_title = "Gamma peak frequency"
slide1_rows = [
    "A_gamma_spectrum_chan3.jpg",
]
add_slide_with_grid(prs, slide1_title, slide1_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 1A: Breaths aligned to breath onset
# ----------------------------------------------------------------
slide1A_title = "Breaths aligned to respiration onset"
slide1A_rows = [
    "allSniffs.jpg",
]
add_slide_with_grid(prs, slide1A_title, slide1A_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 1B: Gamma event counts
# ----------------------------------------------------------------
slide1B_title = "Gamma event counts per condition"
slide1B_rows = [
    "I_gamma_event_counts.jpg",
]
add_slide_with_grid(prs, slide1B_title, slide1B_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 2: ERPs of gamma bursts at peak frequency
# ----------------------------------------------------------------
slide2_title = "ERPs of gamma bursts at peak frequency"
slide2_rows = [
    "B_ERP_sem_chan3.jpg",                     # ERP w/ SEM
    "rawSingleEventTrace_neutral_evt*.jpg",    # random neutral event
    "rawSingleEventTrace_happy_evt*.jpg",      # random happy event
    "rawSingleEventTrace_sad_evt*.jpg",        # random sad event
]
add_slide_with_grid(prs, slide2_title, slide2_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 3: Gamma phase predictability
# ----------------------------------------------------------------
slide3_title = "Gamma phase predictability across cycles"
slide3_rows = [
    "C_gammaPhase_afterKcycles_K1cycles.jpg",
    "C_gammaPhase_afterKcycles_K2cycles.jpg",
    "C_gammaPhase_afterKcycles_K3cycles.jpg",
    "C_gammaPhase_afterKcycles_K4cycles.jpg",
    "C_gammaPhase_afterKcycles_K5cycles.jpg",
]
add_slide_with_grid(prs, slide3_title, slide3_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 4: Gamma bursts are aligned to inhale
# ----------------------------------------------------------------
slide4_title = "Gamma bursts are aligned to inhalation"
slide4_rows = [
    "E_resp_phase_at_gamma_onset.jpg",
    "extra_gamma_power_locked_to_resp_onset.jpg",
]
add_slide_with_grid(prs, slide4_title, slide4_rows, rootPath, subject_order, subj_short)


# ----------------------------------------------------------------
# Slide 5: Gamma burst characteristics
# ----------------------------------------------------------------
slide5_title = "Gamma burst characteristics"
slide5_rows = [
    "extra_gamma_ROI_boxplot_pm100ms_pm5Hz.jpg",
    "extra_IBI_freqpoly_by_condition.jpg",
    "extra_burst_autocorr_pm5s.jpg",
    "D_burst_length_distribution.jpg",
]
add_slide_with_grid(prs, slide5_title, slide5_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 6: Gamma burst relation to other frequencies
# ----------------------------------------------------------------
slide6_title = "Gamma burst relation to other frequencies"
slide6_rows = [
    "extra_delta_power_around_gamma_pm2s.jpg",
    "extra_theta_power_around_gamma_pm2s.jpg",
    "extra_alpha_power_around_gamma_pm2s.jpg",
    "extra_beta_power_around_gamma_pm2s.jpg",
]
add_slide_with_grid(prs, slide6_title, slide6_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 7: Gamma burst relation to other frequencies 2
# ----------------------------------------------------------------
slide7_title = "Gamma burst relation to other frequencies (phase)"
slide7_rows = [
    "F_phase_2Hz_at_onset.jpg",
    "F_phase_5Hz_at_onset.jpg",
    "F_phase_8Hz_at_onset.jpg",
    "F_phase_12Hz_at_onset.jpg",
    "F_phase_15Hz_at_onset.jpg",
]
add_slide_with_grid(prs, slide7_title, slide7_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 8: time frequency power locked to gamma bursts
# ----------------------------------------------------------------
slide8_title = "Time–frequency power locked to gamma bursts"
slide8_rows = [
    "TF_power_neutral_pm1s.jpg",
    "TF_power_happy_pm1s.jpg",
    "TF_power_sad_pm1s.jpg",
]
add_slide_with_grid(prs, slide8_title, slide8_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 9: time frequency power differences
# ----------------------------------------------------------------
slide9_title = "Time–frequency power differences (power)"
slide9_rows = [
    "extra_TF_diff_Happy_minus_Neutral_pm1s.jpg",
    "extra_TF_diff_Sad_minus_Neutral_pm1s.jpg",
]
add_slide_with_grid(prs, slide9_title, slide9_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 10: time frequency ITPC locked to gamma bursts
# ----------------------------------------------------------------
slide10_title = "Time–frequency phase consistency (ITPC) locked to gamma bursts"
slide10_rows = [
    "TF_ITPC_neutral_pm1s.jpg",
    "TF_ITPC_happy_pm1s.jpg",
    "TF_ITPC_sad_pm1s.jpg",
]
add_slide_with_grid(prs, slide10_title, slide10_rows, rootPath, subject_order, subj_short)

# ----------------------------------------------------------------
# Slide 11: time frequency ITPC differences
# ----------------------------------------------------------------
slide11_title = "Time–frequency ITPC differences"
slide11_rows = [
    "extra_TF_ITPC_diff_Happy_minus_Neutral_pm1s.jpg",
    "extra_TF_ITPC_diff_Sad_minus_Neutral_pm1s.jpg",
]
add_slide_with_grid(prs, slide11_title, slide11_rows, rootPath, subject_order, subj_short)
# ============================================================
# SAVE
# ============================================================
prs.save(r"R:\Neurology\Zelano_Lab\Lab_Common\Adam\Dupi_processing\EmotionalMovieTask.pptx")
print(f"Saved slide deck to: {output_pptx}")
