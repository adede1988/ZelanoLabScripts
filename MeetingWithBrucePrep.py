# -*- coding: utf-8 -*-
"""
Created on Tue Jan 13 15:09:09 2026

@author: dtf8829
"""

"""
Auto-build a PowerPoint slide deck with a 3x3 grid per session.

Layout per slide:
  Columns = tasks: breathingTask | cueTask | O15
  Rows    = figures: allSniffs.jpg | macroSpectra.jpg | sniff_TF.jpg

Each cell shows the corresponding image from:
  BASE_DIR\<sessionID>\<taskFolder>\<figureName>

Output:
  BASE_DIR\Dupi_processing_summary.pptx

Dependencies:
  pip install python-pptx pillow
"""

import os
import glob
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

try:
    from PIL import Image
    PIL_OK = True
except Exception:
    PIL_OK = False


# -----------------------------
# User config
# -----------------------------
BASE_DIR = Path(r"R:\Neurology\Zelano_Lab\Lab_Common\Adam\Dupi_processing")

SESSION_IDS = [
    "250818_Dupi_NMH_JH_1",
    "250623_DUPI_NMH_KS_2",
    "250623_Dupi_NMH_KS_1",
    "250818_Dupi_NMH_JH_2",
    "250811_Dupi_NMH_TPB_1",
    "250811_Dupi_NMH_TB_2",
    "250929_Dupi_NMH_GH_1",
    "251002_Dupi_NMH_AB_1",
    "251027_Dupi_NMH_DL_1",
    "250929_Dupi_NMH_GH_2",
    "251002_Dupi_NMH_AB_2",
    "251013_Dupi_NMH_JN_2",
]

TASK_FOLDERS = ["breathingTask", "cueTask", "O15"]
FIG_FILES = ["allSniffs.jpg", "macroSpectra.jpg", "sniff_TF.jpg"]  # rows in this order

OUT_PPTX = BASE_DIR / "Dupi_processing_summary.pptx"


# -----------------------------
# Helpers
# -----------------------------
def find_case_insensitive_dir(parent: Path, name: str) -> Path | None:
    """Return Path to a child directory matching name case-insensitively."""
    target = name.lower()
    if (parent / name).is_dir():
        return parent / name
    for p in parent.iterdir():
        if p.is_dir() and p.name.lower() == target:
            return p
    return None


def find_image_file(folder: Path, fname: str) -> Path | None:
    """Find an image file in folder. Try exact match, then case-insensitive glob."""
    p = folder / fname
    if p.exists():
        return p

    # Case-insensitive search by basename
    base = Path(fname).stem.lower()
    ext = Path(fname).suffix.lower()
    candidates = []
    for q in folder.glob("*"):
        if q.is_file():
            if q.suffix.lower() == ext and q.stem.lower() == base:
                candidates.append(q)
    if candidates:
        return candidates[0]

    # Fallback: contains match (helps if filenames vary slightly)
    patt = str(folder / f"*{base}*{ext}")
    hits = glob.glob(patt)
    return Path(hits[0]) if hits else None


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=(0, 0, 0), align_center=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    if align_center:
        p.alignment = 1  # PP_ALIGN.CENTER
    return tb


def add_image_fit(slide, img_path: Path, box_left, box_top, box_w, box_h):
    """
    Add image fitted into the box while preserving aspect ratio (requires PIL).
    If PIL missing, image will be stretched to the box.
    """
    if not img_path or not img_path.exists():
        # placeholder
        shp = slide.shapes.add_shape(
            1, box_left, box_top, box_w, box_h  # MSO_SHAPE.RECTANGLE = 1
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(245, 245, 245)
        shp.line.color.rgb = RGBColor(180, 180, 180)
        add_textbox(slide, box_left, box_top, box_w, box_h, "MISSING", font_size=14, bold=True,
                    color=(140, 0, 0), align_center=True)
        return

    if not PIL_OK:
        slide.shapes.add_picture(str(img_path), box_left, box_top, width=box_w, height=box_h)
        return

    with Image.open(img_path) as im:
        iw, ih = im.size

    # Fit image in box preserving aspect
    box_w_in = box_w / 914400  # EMU per inch
    box_h_in = box_h / 914400

    img_ar = iw / ih
    box_ar = box_w_in / box_h_in

    if img_ar >= box_ar:
        # limited by width
        w_in = box_w_in
        h_in = box_w_in / img_ar
    else:
        # limited by height
        h_in = box_h_in
        w_in = box_h_in * img_ar

    left = box_left + Inches((box_w_in - w_in) / 2)
    top = box_top + Inches((box_h_in - h_in) / 2)

    slide.shapes.add_picture(str(img_path), left, top, width=Inches(w_in), height=Inches(h_in))


def set_widescreen(prs: Presentation):
    # 13.333 x 7.5 inches (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


# -----------------------------
# Main build
# -----------------------------
def build_deck():
    if not BASE_DIR.exists():
        raise FileNotFoundError(f"BASE_DIR does not exist: {BASE_DIR}")

    prs = Presentation()
    set_widescreen(prs)
    blank = prs.slide_layouts[6]

    # Layout constants
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin_l = Inches(0.35)
    margin_r = Inches(0.25)
    margin_b = Inches(0.30)
    title_h  = Inches(0.45)
    header_h = Inches(0.35)

    row_label_w = Inches(1.35)
    grid_left = margin_l + row_label_w
    grid_top  = margin_l + title_h + header_h
    grid_w = slide_w - grid_left - margin_r
    grid_h = slide_h - grid_top - margin_b

    cell_w = grid_w / 3
    cell_h = grid_h / 3

    # -----------------------------
    # Spike artifact pipeline slides (250818_Dupi_NMH_JH_1)
    # Drop this block inside build_deck(), AFTER your session slides loop, BEFORE prs.save(...)
    # Requires helper functions already in your script: add_textbox(), add_image_fit(), find_image_file()
    # -----------------------------
    
    
    def add_bullets(slide, left, top, width, height, bullets, font_size=14):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(font_size)
        return tb
    
    # Folder with the example pipeline figures
    spike_demo_sid = "250818_Dupi_NMH_JH_1"
    spike_dir = BASE_DIR / spike_demo_sid
    
    # Layout numbers (inches) matching your deck (16:9 widescreen)
    slide_w_in = slide_w / 914400
    slide_h_in = slide_h / 914400
    ml, mr, mb = 0.35, 0.25, 0.30
    title_y = 0.10
    title_h = 0.45
    content_top = 0.75
    gap = 0.25
    
    usable_w_in = slide_w_in - ml - mr
    usable_h_in = slide_h_in - content_top - mb
    
    img_w_in = usable_w_in * 0.62
    txt_w_in = usable_w_in - img_w_in - gap
    
    img_left = Inches(ml)
    txt_left = Inches(ml + img_w_in + gap)
    content_top_len = Inches(content_top)
    
    # Slide specs: (title, [images], bullets, optional image captions)
    slides_spec = [
        (
            "1) Raw intranasal macro data: transient amplifier spikes",
            ["macrosRaw_zoom.jpg"],
            [
                "Raw intranasal macro channels show transient amplifier spike artifacts.",
                "We are using the hardware in a non-standard configuration, so the exact source isn’t fully clear.",
                "Importantly, these intranasal macro spikes look qualitatively different from typical scalp EEG noise."
            ],
            []
        ),
        (
            "2) Bipolar referencing reduces shared noise",
            ["bipolarAll_zoom.jpg"],
            [
                "Bipolar referencing removes a large fraction of shared/common noise.",
                "Implemented by differencing adjacent contacts: (2−1), (3−2), (4−3), …",
                "This improves interpretability and prepares the data for spike-focused processing."
            ],
            []
        ),
        (
            "3) Split into high- and low-frequency components",
            ["bipolarHigh_zoom.jpg", "bipolarLow_zoom.jpg"],
            [
                "Spike artifacts are predominantly high-frequency.",
                "We split the bipolar signal into: high (>10 Hz) and low (≤10 Hz) components.",
                "Spikes become easier to detect and threshold in the high-frequency component."
            ],
            ["High (>10 Hz)", "Low (≤10 Hz)"]
        ),
        (
            "4) Spike detection features in a sliding window",
            ["spikeDetection_zoom.jpg"],
            [
                "Within a 22 ms sliding window we compute peak-to-peak range (max − min).",
                "We also compute peak prominence (rise relative to shoulders).",
                "These features provide a robust basis for spike candidate detection."
            ],
            []
        ),
        (
            "5) ICA-based spike component identification and surgical removal",
            ["ICA_fig.jpg"],
            [
                "High-frequency bipolar data are submitted to ICA.",
                "We compare component variability at spike-detection times vs non-spike times to automatically select the spike IC.",
                "Data are reconstructed with the spike IC surgically windowed out at detected spike times."
            ],
            []
        ),
        (
            "6) Cleaned high-frequency bipolar signal",
            ["spikeClean_zoom.jpg"],
            [
                "After ICA-based removal, the high-frequency component is substantially cleaner.",
                "Residual spike transients are strongly reduced while preserving physiological high-frequency activity."
            ],
            []
        ),
        (
            "7) Final cleaned signal: add low-frequency component back",
            ["spikeCleanFull_zoom.jpg"],
            [
                "Finally, the low-frequency (≤10 Hz) component is added back to the cleaned high-frequency signal.",
                "This yields the final cleaned broadband signal used for downstream analyses."
            ],
            []
        ),
    ]
    
    for title, imgs, bullets, captions in slides_spec:
        slide = prs.slides.add_slide(blank)
    
        # Title
        add_textbox(
            slide,
            Inches(ml), Inches(title_y),
            Inches(usable_w_in), Inches(title_h),
            text=title, font_size=20, bold=True, align_center=False
        )
    
        # Text bullets on right
        add_bullets(
            slide,
            txt_left, content_top_len,
            Inches(txt_w_in), Inches(usable_h_in),
            bullets=bullets,
            font_size=14
        )
    
        # Images on left
        if len(imgs) == 1:
            img_path = find_image_file(spike_dir, imgs[0]) if spike_dir.exists() else None
            add_image_fit(
                slide,
                img_path,
                img_left, content_top_len,
                Inches(img_w_in), Inches(usable_h_in)
            )
        else:
            # stack two images vertically in the left panel (for Slide 3)
            half_h_in = (usable_h_in - 0.15) / 2
            img1_top = content_top
            img2_top = content_top + half_h_in + 0.15
    
            img_path1 = find_image_file(spike_dir, imgs[0]) if spike_dir.exists() else None
            img_path2 = find_image_file(spike_dir, imgs[1]) if spike_dir.exists() else None
    
            add_image_fit(slide, img_path1, img_left, Inches(img1_top), Inches(img_w_in), Inches(half_h_in))
            add_image_fit(slide, img_path2, img_left, Inches(img2_top), Inches(img_w_in), Inches(half_h_in))
    
            # Optional captions above each stacked image
            if len(captions) == 2:
                cap_h = 0.22
                add_textbox(slide, img_left, Inches(img1_top - cap_h), Inches(img_w_in), Inches(cap_h),
                            text=captions[0], font_size=12, bold=True, align_center=True)
                add_textbox(slide, img_left, Inches(img2_top - cap_h), Inches(img_w_in), Inches(cap_h),
                            text=captions[1], font_size=12, bold=True, align_center=True)
    


    # -----------------------------
    # Add group-level summary slide (two figures side-by-side)
    # Place this block AFTER the session slides loop, BEFORE prs.save(...)
    # -----------------------------
    group_dir = BASE_DIR / "groupStatFigs"
    img_acc = group_dir / "performance_accuracy_session1to2.png"
    img_dp  = group_dir / "performance_dprime_session1to2.png"
    
    slide = prs.slides.add_slide(blank)
    
    # Title
    add_textbox(
        slide,
        margin_l, Inches(0.10),
        slide_w - margin_l - margin_r, Inches(0.45),
        text="Group performance: Session 1 → 2",
        font_size=22, bold=True, align_center=False
    )
    
    # Two-column layout
    top = Inches(0.75)
    usable_h = slide_h - top - margin_b
    gap = Inches(0.20)
    
    col_w = (slide_w - margin_l - margin_r - gap) / 2
    left_x = margin_l
    right_x = margin_l + col_w + gap
    
    # Optional small panel titles
    add_textbox(slide, left_x, top - Inches(0.30), col_w, Inches(0.25),
                text="Accuracy", font_size=14, bold=True, align_center=True)
    add_textbox(slide, right_x, top - Inches(0.30), col_w, Inches(0.25),
                text="d′", font_size=14, bold=True, align_center=True)
    
    # Images (fit preserving aspect ratio if PIL is available; otherwise stretch)
    add_image_fit(slide, img_acc, left_x,  top, col_w, usable_h)
    add_image_fit(slide, img_dp,  right_x, top, col_w, usable_h)       


    # -----------------------------
    # Add group-level slide: lowGammaPeakAcrossResp.png (single figure)
    # Place this block AFTER the session slides loop, BEFORE prs.save(...)
    # -----------------------------
    group_dir = BASE_DIR / "groupStatFigs"
    img_lowgamma = group_dir / "lowGammaPeakAcrossResp.png"
    
    slide = prs.slides.add_slide(blank)
    
    # Title
    add_textbox(
        slide,
        margin_l, Inches(0.10),
        slide_w - margin_l - margin_r, Inches(0.45),
        text="Low-gamma peak across respiration",
        font_size=22, bold=True, align_center=False
    )
    
    # Image full-width (with margins)
    top = Inches(0.70)
    usable_h = slide_h - top - margin_b
    usable_w = slide_w - margin_l - margin_r
    
    add_image_fit(slide, img_lowgamma, margin_l, top, usable_w, usable_h)

    # Row/col labels
    row_labels = ["allSniffs", "macroSpectra", "sniff_TF"]
    col_labels = TASK_FOLDERS

    # Precompute available session folder names for case-insensitive mapping
    session_dirs = {p.name.lower(): p for p in BASE_DIR.iterdir() if p.is_dir()}

    for sid in SESSION_IDS:
        slide = prs.slides.add_slide(blank)

        # Title
        add_textbox(
            slide,
            margin_l, Inches(0.10),
            slide_w - margin_l - margin_r, title_h,
            text=sid,
            font_size=20, bold=True, align_center=False
        )

        # Column headers
        for j, col in enumerate(col_labels):
            left = grid_left + j * cell_w
            add_textbox(
                slide,
                left, margin_l + title_h,
                cell_w, header_h,
                text=col,
                font_size=14, bold=True, align_center=True
            )

        # Row labels
        for i, rlab in enumerate(row_labels):
            top = grid_top + i * cell_h
            add_textbox(
                slide,
                margin_l, top,
                row_label_w, cell_h,
                text=rlab,
                font_size=14, bold=True, align_center=True
            )

        # Resolve session folder
        sess_dir = None
        if (BASE_DIR / sid).is_dir():
            sess_dir = BASE_DIR / sid
        else:
            sess_dir = session_dirs.get(sid.lower(), None)

        if sess_dir is None:
            # Fill whole grid with "missing session folder"
            add_textbox(slide, grid_left, grid_top, grid_w, grid_h,
                        f"Missing session folder:\n{sid}", font_size=18, bold=True,
                        color=(140, 0, 0), align_center=True)
            continue

        # Add images
        for j, task in enumerate(TASK_FOLDERS):
            task_dir = find_case_insensitive_dir(sess_dir, task)
            for i, fig in enumerate(FIG_FILES):
                box_left = grid_left + j * cell_w
                box_top  = grid_top + i * cell_h

                # Padding inside each cell
                pad = Inches(0.08)
                inner_left = box_left + pad
                inner_top  = box_top + pad
                inner_w = cell_w - 2 * pad
                inner_h = cell_h - 2 * pad

                img_path = None
                if task_dir is not None:
                    img_path = find_image_file(task_dir, fig)

                add_image_fit(slide, img_path, inner_left, inner_top, inner_w, inner_h)
                
         

    prs.save(str(OUT_PPTX))
    print(f"Saved: {OUT_PPTX}")


if __name__ == "__main__":
    build_deck()
